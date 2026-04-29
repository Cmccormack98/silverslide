"""
python-pptx slide builder for SilverSlide Agent.
Default builder — no Node.js required.

Senior-friendly design:
  • Title/summary slides: deep navy background, white text
  • Content slides: warm off-white background, navy title bar
  • Video slide: light blue tint, thumbnail + clickable link
  • Fonts: Calibri throughout  (clear, universally available)
  • Title bar:  30 pt bold   |  Body bullets: 24 pt
  • Left amber stripe for visual warmth on every slide
"""

import io
import os
import base64
from pathlib import Path
from typing import Optional

import requests as _http

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree


# ── Palette ──────────────────────────────────────────────────────────────────
class C:
    navy       = RGBColor(0x1B, 0x3A, 0x5C)
    navyDeep   = RGBColor(0x12, 0x27, 0x43)
    amber      = RGBColor(0xE0, 0x7B, 0x39)
    white      = RGBColor(0xFF, 0xFF, 0xFF)
    offWhite   = RGBColor(0xF8, 0xF6, 0xF2)
    textDark   = RGBColor(0x1A, 0x20, 0x2C)
    steelBlue  = RGBColor(0x9E, 0xC8, 0xE0)
    slideNum   = RGBColor(0xA0, 0xAE, 0xC0)
    videoBg    = RGBColor(0xEE, 0xF6, 0xFB)
    linkBlue   = RGBColor(0x25, 0x63, 0xEB)
    captionGray= RGBColor(0x6B, 0x8F, 0xAF)


FONT = "Calibri"


def _fetch_image(hint: str) -> Optional[io.BytesIO]:
    """Fetch a relevant landscape photo from Pexels. Returns None if unavailable."""
    key = os.environ.get("PEXELS_API_KEY", "")   # read fresh every call
    if not key or not hint:
        return None
    try:
        search = _http.get(
            "https://api.pexels.com/v1/search",
            params={"query": hint, "per_page": 1, "orientation": "landscape"},
            headers={"Authorization": key},
            timeout=6,
        )
        if search.status_code == 200:
            photos = search.json().get("photos", [])
            if photos:
                img_url = photos[0]["src"]["large"]
                img_resp = _http.get(img_url, timeout=8)
                if img_resp.status_code == 200:
                    return io.BytesIO(img_resp.content)
    except Exception:
        pass
    return None

# Slide canvas: 10" × 5.625"  (LAYOUT_16x9)
W = Inches(10)
H = Inches(5.625)


# ── Low-level helpers ─────────────────────────────────────────────────────────

def _add_rect(slide, x, y, w, h, color: RGBColor, *, border=False):
    """Add a filled rectangle with no visible border (unless border=True)."""
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border:
        shape.line.color.rgb = color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()  # no line
    return shape


def _add_text(
    slide,
    text: str,
    x, y, w, h,
    size: int,
    color: RGBColor,
    *,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    valign="top",
    font=FONT,
    underline=False,
    url: Optional[str] = None,
) -> None:
    """Add a simple single-run text box."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.text_frame.word_wrap = True

    if valign == "middle":
        from pptx.enum.text import MSO_ANCHOR
        txBox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = txBox.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    run.font.underline = underline
    if url:
        run.hyperlink.address = url


def _set_bg(slide, color: RGBColor):
    """Set slide background to a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_notes(slide, text: str):
    """Set speaker notes on a slide."""
    if not text:
        return
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


def _add_title_bar(slide, title: str):
    """Navy bar across the top + white title text (used on content slides)."""
    _add_rect(slide, 0, 0, 10, 1.15, C.navy)
    _add_text(
        slide, title,
        x=0.45, y=0.12, w=9.1, h=0.93,
        size=30, color=C.white, bold=True,
        align=PP_ALIGN.LEFT, valign="middle",
    )


def _add_left_stripe(slide, y_start=0.0, height=5.625):
    """Amber left stripe — adds warmth and visual structure."""
    _add_rect(slide, 0, y_start, 0.28, height, C.amber)


def _add_slide_num(slide, current: int, total: int):
    _add_text(
        slide, f"{current} / {total}",
        x=8.9, y=5.25, w=0.9, h=0.3,
        size=11, color=C.slideNum, align=PP_ALIGN.RIGHT,
    )


def _add_bullets(slide, bullets: list[dict], x, y, w, h):
    """
    Add a bulleted list.  Each dict has key 'text'.
    Uses 24 pt, dark text, with spacing between bullets.
    """
    if not bullets:
        return
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT

        # Bullet character via paragraph XML
        pPr = p._pPr
        if pPr is None:
            pPr = etree.SubElement(p._p, qn("a:pPr"))
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", "\u2022")  # •

        # Space between bullets
        spcAft = etree.SubElement(pPr, qn("a:spcAft"))
        spcPts = etree.SubElement(spcAft, qn("a:spcPts"))
        spcPts.set("val", "1000")  # 10 pt after paragraph

        run = p.add_run()
        run.text = bullet["text"]
        run.font.size = Pt(24)
        run.font.color.rgb = C.textDark
        run.font.name = FONT


# ── Builder ───────────────────────────────────────────────────────────────────

def build_pptx(payload: dict, output_path: str) -> None:
    """
    Build a senior-friendly .pptx from the payload dict and save to output_path.

    payload keys:
        title        str
        slides       list of slide dicts (slide_type, title, bullets, speaker_notes)
        video        dict or None (video_id, title, channel, duration, url, thumbnail_base64)
    """
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    # Blank layout (index 6) — we place everything manually
    blank = prs.slide_layouts[6]

    slides_data   = payload["slides"]
    title_slide   = next((s for s in slides_data if s["slide_type"] == "title"),   None)
    content_slides= [s for s in slides_data if s["slide_type"] == "content"]
    summary_slide = next((s for s in slides_data if s["slide_type"] == "summary"), None)
    video         = payload.get("video")
    has_video     = bool(video and video.get("url"))

    total = 1 + len(content_slides) + (1 if has_video else 0) + 1
    cur   = 0

    # ══════════════════════════════════════════════════════════════════════
    # SLIDE 1 — TITLE
    # ══════════════════════════════════════════════════════════════════════
    cur += 1
    s1 = prs.slides.add_slide(blank)
    _set_bg(s1, C.navy)
    _add_left_stripe(s1, 0, 5.625)

    # Bottom footer bar
    _add_rect(s1, 0, 5.125, 10, 0.5, C.navyDeep)

    # Amber divider line
    _add_rect(s1, 0.55, 2.9, 9.0, 0.04, C.amber)

    # Main title
    _add_text(
        s1, payload["title"],
        x=0.55, y=0.85, w=9.0, h=1.9,
        size=40, color=C.white, bold=True,
        align=PP_ALIGN.LEFT,
    )

    # Tagline / subtitle
    tagline = title_slide["title"] if title_slide and title_slide["title"] != payload["title"] else ""
    if tagline:
        _add_text(
            s1, tagline,
            x=0.55, y=3.05, w=9.0, h=0.75,
            size=20, color=C.steelBlue,
            align=PP_ALIGN.LEFT,
        )

    _add_text(
        s1,
        "Created by SilverSlide Agent  \u2022  Senior-Friendly Presentation",
        x=0.55, y=5.17, w=8.5, h=0.35,
        size=11, color=C.captionGray,
        align=PP_ALIGN.LEFT,
    )

    if title_slide and title_slide.get("speaker_notes"):
        _add_notes(s1, title_slide["speaker_notes"])

    # ══════════════════════════════════════════════════════════════════════
    # CONTENT SLIDES
    # ══════════════════════════════════════════════════════════════════════
    for slide_data in content_slides:
        cur += 1
        sl = prs.slides.add_slide(blank)
        _set_bg(sl, C.offWhite)
        _add_title_bar(sl, slide_data["title"])
        _add_rect(sl, 0, 1.15, 0.28, 5.625 - 1.15, C.amber)  # amber stripe below bar

        # Try to fetch a relevant image
        hint      = (slide_data.get("image_hint") or slide_data["title"])
        img_stream = _fetch_image(hint)

        if img_stream:
            # Narrower bullet area on left, image panel on right
            if slide_data.get("bullets"):
                _add_bullets(sl, slide_data["bullets"], x=0.5, y=1.28, w=5.1, h=4.0)
            # Subtle navy frame behind image
            _add_rect(sl, 5.82, 1.18, 4.05, 4.27, C.navy)
            try:
                sl.shapes.add_picture(img_stream, Inches(5.88), Inches(1.24), Inches(3.93), Inches(4.15))
            except Exception:
                # Image failed — fall back to full-width
                if slide_data.get("bullets"):
                    _add_bullets(sl, slide_data["bullets"], x=0.5, y=1.28, w=9.2, h=4.0)
        else:
            if slide_data.get("bullets"):
                _add_bullets(sl, slide_data["bullets"], x=0.5, y=1.28, w=9.2, h=4.0)

        _add_slide_num(sl, cur, total)

        if slide_data.get("speaker_notes"):
            _add_notes(sl, slide_data["speaker_notes"])

    # ══════════════════════════════════════════════════════════════════════
    # VIDEO SLIDE
    # ══════════════════════════════════════════════════════════════════════
    if has_video:
        cur += 1
        vs = prs.slides.add_slide(blank)
        _set_bg(vs, C.videoBg)
        _add_title_bar(vs, "Watch This Short Video")
        _add_rect(vs, 0, 1.15, 0.28, 5.625 - 1.15, C.amber)

        thumb_x, thumb_y, thumb_w, thumb_h = 0.45, 1.3, 5.4, 3.05

        # Thumbnail image
        thumb_b64 = video.get("thumbnail_base64")
        if thumb_b64:
            try:
                # Strip the data-URI prefix: "image/jpeg;base64,..."
                raw = thumb_b64.split(",", 1)[-1]
                img_bytes = base64.b64decode(raw)
                img_stream = io.BytesIO(img_bytes)
                vs.shapes.add_picture(
                    img_stream,
                    Inches(thumb_x), Inches(thumb_y),
                    Inches(thumb_w), Inches(thumb_h),
                )
            except Exception:
                # Grey placeholder on error
                _add_rect(vs, thumb_x, thumb_y, thumb_w, thumb_h, RGBColor(0xCC, 0xCC, 0xCC))
        else:
            _add_rect(vs, thumb_x, thumb_y, thumb_w, thumb_h, RGBColor(0xCC, 0xCC, 0xCC))

        # Dark info panel (right side)
        _add_rect(vs, 6.1, 1.3, 3.65, 3.05, C.navy)

        _add_text(vs, "HOW TO PLAY",
                  x=6.2, y=1.42, w=3.45, h=0.35,
                  size=11, color=C.amber, bold=True, align=PP_ALIGN.CENTER)

        _add_text(vs, "Click the link below\nto open the video",
                  x=6.2, y=1.82, w=3.45, h=0.75,
                  size=15, color=C.white, bold=True, align=PP_ALIGN.CENTER)

        _add_rect(vs, 6.35, 2.65, 3.1, 0.03, RGBColor(0x3A, 0x5A, 0x7A))

        # Video title
        title_text = (video.get("title") or "")[:80]
        _add_text(vs, title_text,
                  x=6.2, y=2.72, w=3.45, h=0.85,
                  size=13, color=C.steelBlue, align=PP_ALIGN.CENTER)

        if video.get("duration"):
            _add_text(vs, f"Duration: {video['duration']}",
                      x=6.2, y=3.62, w=3.45, h=0.3,
                      size=12, color=RGBColor(0x7A, 0xAC, 0xCB), align=PP_ALIGN.CENTER)

        if video.get("channel"):
            _add_text(vs, f"From: {video['channel']}",
                      x=6.2, y=3.95, w=3.45, h=0.28,
                      size=11, color=RGBColor(0x5A, 0x8A, 0xAA), align=PP_ALIGN.CENTER)

        # Clickable text link below thumbnail
        _add_text(
            vs,
            "\u25b6  Click here to watch this video",
            x=thumb_x, y=thumb_y + thumb_h + 0.12,
            w=thumb_w, h=0.38,
            size=16, color=C.linkBlue,
            align=PP_ALIGN.CENTER, bold=True, underline=True,
            url=video["url"],
        )

        _add_slide_num(vs, cur, total)
        _add_notes(vs, f"Video: {video.get('title', '')}\n"
                       f"Channel: {video.get('channel', '')}\n"
                       f"Duration: {video.get('duration', '')}\n"
                       f"URL: {video.get('url', '')}")

    # ══════════════════════════════════════════════════════════════════════
    # SUMMARY SLIDE
    # ══════════════════════════════════════════════════════════════════════
    cur += 1
    ss = prs.slides.add_slide(blank)
    _set_bg(ss, C.navy)
    _add_left_stripe(ss, 0, 5.625)
    _add_rect(ss, 0, 5.125, 10, 0.5, C.navyDeep)

    summary_title = (summary_slide["title"] if summary_slide and summary_slide["title"]
                     else "What to Remember")
    _add_text(
        ss, summary_title,
        x=0.5, y=0.2, w=9.2, h=0.9,
        size=34, color=C.white, bold=True,
        align=PP_ALIGN.LEFT, valign="middle",
    )

    # Amber divider
    _add_rect(ss, 0.5, 1.18, 9.2, 0.04, C.amber)

    # Checkmark bullets (plain text with ✓ prefix — avoids double-bullet issue)
    if summary_slide and summary_slide.get("bullets"):
        txBox = ss.shapes.add_textbox(
            Inches(0.5), Inches(1.3), Inches(9.2), Inches(3.5)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(summary_slide["bullets"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT

            # Space after paragraph
            pPr = p._pPr
            if pPr is None:
                pPr = etree.SubElement(p._p, qn("a:pPr"))
            spcAft = etree.SubElement(pPr, qn("a:spcAft"))
            spcPts = etree.SubElement(spcAft, qn("a:spcPts"))
            spcPts.set("val", "1400")  # 14 pt after

            run = p.add_run()
            run.text = f"\u2713  {bullet['text']}"
            run.font.size = Pt(24)
            run.font.color.rgb = C.white
            run.font.name = FONT

    _add_text(
        ss,
        "Questions? Talk with your doctor, family, or community coordinator.",
        x=0.5, y=5.14, w=9.0, h=0.35,
        size=12, color=C.captionGray, italic=True,
        align=PP_ALIGN.LEFT,
    )

    if summary_slide and summary_slide.get("speaker_notes"):
        _add_notes(ss, summary_slide["speaker_notes"])

    prs.save(output_path)
